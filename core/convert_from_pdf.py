"""PDF -> Word / Excel / PowerPoint / Images. All pure Python, no Office needed."""
import re
import tempfile
from pathlib import Path

import fitz
from openpyxl import Workbook
from pdf2docx import Converter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

from core import legibility, ocr

EMU_PER_POINT = 12700  # exact: 914400 EMU/inch / 72 points/inch
PDF_BOLD_FLAG = 1 << 4
PDF_ITALIC_FLAG = 1 << 1
PUA_START, PUA_END = 0xE000, 0xF8FF  # Unicode Private Use Area


def _sanitize_text(text: str) -> str:
    """Bullet-list markers in PowerPoint-exported PDFs are usually a
    character from a symbol font (Wingdings, Webdings, etc.) rendered via
    that font's own private glyph mapping, e.g. Wingdings code point
    U+F0D8 is a right-pointing arrow. Extracted as plain text and dropped
    into a normal font, those code points have no defined glyph at all and
    render as broken boxes throughout the deck (this was the "huge issue").
    Since we don't carry the original font family over (see _add_text_block),
    there's no font install that would make the real glyph show up either,
    swapping in a plain bullet character preserves what the mark actually
    meant instead of showing a hole in the page."""
    return "".join("•" if PUA_START <= ord(c) <= PUA_END else c for c in text)


def pdf_to_word(path: str, save_path: str) -> None:
    cv = Converter(path)
    try:
        cv.convert(save_path)
    finally:
        cv.close()


def _cluster_1d(values: list[float], tol: float) -> list[float]:
    """Collapse near-equal scalars into their cluster means. Used to turn the
    many tiny ruling-line segments a PDF draws for one visual gridline into a
    single representative coordinate."""
    values = sorted(values)
    if not values:
        return []
    groups = [[values[0]]]
    for v in values[1:]:
        if v - groups[-1][-1] <= tol:
            groups[-1].append(v)
        else:
            groups.append([v])
    return [sum(g) / len(g) for g in groups]


def _ruled_table(page):
    """Best case: the table is drawn with real gridlines (most bank
    statements are). Cluster every vertical and horizontal edge into a set of
    explicit column/row lines and let pdfplumber cut the page on exactly those
    lines. This groups a transaction whose text wraps onto three physical
    lines back into the single bordered cell it visually occupies, which the
    default extract_tables() cannot do (it instead reports each bordered row
    as its own separate tiny table). Returns None when there isn't a real
    grid to use."""
    vedges = [e["x0"] for e in page.edges if e["orientation"] == "v"]
    hedges = [e["top"] for e in page.edges if e["orientation"] == "h"]
    vlines = _cluster_1d(vedges, tol=3)
    hlines = _cluster_1d(hedges, tol=3)
    if len(vlines) < 3 or len(hlines) < 3:
        return None
    settings = {
        "vertical_strategy": "explicit",
        "horizontal_strategy": "explicit",
        "explicit_vertical_lines": vlines,
        "explicit_horizontal_lines": hlines,
        "text_x_tolerance": 1,
        "text_y_tolerance": 1,
    }
    table = page.extract_table(settings) or []
    rows = [r for r in table if any((c or "").strip() for c in r)]
    return rows or None


def _column_bounds(words: list[dict], page_width: float) -> list[float]:
    """For a borderless table (no gridlines), infer column separators from the
    vertical whitespace gaps that persist down the page. Mark every x the
    words actually cover, then treat each wide empty band between covered
    regions as a column boundary."""
    if not words:
        return [0.0, page_width]
    width = int(page_width) + 2
    covered = bytearray(width)
    heights = []
    for w in words:
        x0 = max(0, int(w["x0"]))
        x1 = min(width - 1, int(w["x1"]) + 1)
        for x in range(x0, x1):
            covered[x] = 1
        heights.append(w["bottom"] - w["top"])
    # A gap must be wider than roughly one character to count as a real column
    # separator, not just the normal space between words in a sentence.
    min_gap = max(6, (sum(heights) / len(heights)) * 0.9)
    bounds = [0.0]
    run_start = None
    for x in range(width):
        if covered[x] == 0:
            if run_start is None:
                run_start = x
        else:
            if run_start is not None and (x - run_start) >= min_gap:
                bounds.append((run_start + x) / 2)
            run_start = None
    bounds.append(float(page_width))
    return bounds


def _grid_from_words(words: list[dict], page_width: float) -> list[list[str]]:
    """Rebuild a table from loose positioned words (used for borderless text
    tables and for OCR output). Words are bucketed into columns by the
    detected boundaries and into rows by vertical position; a line that only
    fills the widest (description) column and leaves the leading columns empty
    is treated as a wrapped continuation and folded back into the row above."""
    if not words:
        return []
    bounds = _column_bounds(words, page_width)
    ncols = len(bounds) - 1

    def col_of(x: float) -> int:
        for i in range(ncols):
            if bounds[i] - 0.5 <= x < bounds[i + 1] - 0.5:
                return i
        return ncols - 1

    # Sort by vertical centre only (never fall through to comparing the word
    # dicts themselves, which raises when two words share the same centre).
    ordered = sorted(words, key=lambda w: (w["top"] + w["bottom"]) / 2)
    med_h = sorted(w["bottom"] - w["top"] for w in words)[len(words) // 2]
    row_tol = max(2.0, med_h * 0.6)

    lines: list[dict] = []
    for w in ordered:
        yc = (w["top"] + w["bottom"]) / 2
        if lines and yc - lines[-1]["yc"] <= row_tol:
            line = lines[-1]
        else:
            line = {"yc": yc, "cells": [[] for _ in range(ncols)]}
            lines.append(line)
        line["yc"] = yc
        line["cells"][col_of((w["x0"] + w["x1"]) / 2)].append((w["x0"], w["text"]))

    # widest column by average text length is the free-text (description) one
    col_len = [0.0] * ncols
    col_n = [0] * ncols
    for ln in lines:
        for i, cell in enumerate(ln["cells"]):
            if cell:
                col_len[i] += sum(len(t) for _, t in cell)
                col_n[i] += 1
    desc_col = max(range(ncols), key=lambda i: (col_len[i] / col_n[i]) if col_n[i] else 0)

    def render(line) -> list[str]:
        out = []
        for cell in line["cells"]:
            cell.sort()
            out.append(" ".join(t for _, t in cell))
        return out

    rows: list[list[str]] = []
    for ln in lines:
        cells = render(ln)
        leading_empty = all(not cells[i].strip() for i in range(ncols) if i != desc_col)
        if rows and leading_empty and cells[desc_col].strip():
            prev = rows[-1]
            prev[desc_col] = (prev[desc_col] + " " + cells[desc_col]).strip()
        elif any(c.strip() for c in cells):
            rows.append(cells)
    return rows


def _winrt_iter(vec):
    """winocr returns Windows runtime vectors (lines/words). Most builds make
    them directly iterable, but fall back to index access so a projection
    quirk can't turn OCR output into silently-empty results."""
    if vec is None:
        return []
    try:
        return list(vec)
    except TypeError:
        try:
            return [vec.get_at(i) for i in range(vec.size)]
        except Exception:
            return []


def _split_line_cols(text: str) -> list[str]:
    """Rough column split for the raw-OCR fallback: OCR tends to preserve the
    gaps between columns as runs of two or more spaces, so split on those."""
    parts = [p for p in re.split(r"\s{2,}", text.strip()) if p]
    return parts or [text.strip()]


def _trim_empty_columns(rows: list[list[str]]) -> list[list[str]]:
    """Drop columns that are empty across every row (the word-grid often leaves
    a blank leading column, and OCR can leave stray empty ones). Rows are first
    padded to equal width so column indexing is safe."""
    if not rows:
        return rows
    ncols = max(len(r) for r in rows)
    norm = [list(r) + [""] * (ncols - len(r)) for r in rows]
    keep = [j for j in range(ncols) if any((norm[i][j] or "").strip() for i in range(len(norm)))]
    if not keep:
        return rows
    return [[r[j] for j in keep] for r in norm]


# A leading date like 06-01, 6/1, or 22/04/2026 at the very start of a cell.
_LEADING_DATE = re.compile(r"^\s*(\d{1,2}[-/]\d{1,2}(?:[-/]\d{2,4})?)\s+(.+)$")


def _split_leading_date(rows: list[list[str]]) -> list[list[str]]:
    """On scanned statements the date and the description often land in one
    column because there's no clean whitespace gap between them. When most rows
    start their first cell with a date, peel that date into its own leading
    column so Date and Description separate. Left untouched when it doesn't
    clearly apply, so it can't scramble a non-statement layout."""
    if not rows:
        return rows
    hits = sum(1 for r in rows if r and _LEADING_DATE.match(r[0] or ""))
    if hits < max(3, 0.3 * len(rows)):
        return rows
    out = []
    for r in rows:
        m = _LEADING_DATE.match(r[0] or "") if r else None
        if m:
            out.append([m.group(1), m.group(2)] + list(r[1:]))
        else:
            out.append([""] + list(r))  # keep columns aligned with the split rows
    return out


def _winocr_page(doc, page_index: int, dpi: int, scale: float):
    """Legacy fallback: Windows' built-in OCR (winocr). Kept only for the case
    where Tesseract isn't reachable; in the frozen exe this often failed to
    import at all, which is why Tesseract is now the primary engine."""
    if not legibility.is_available():
        return None, None
    import io

    import winocr
    from PIL import Image

    page = doc[page_index]
    pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
    result = winocr.recognize_pil_sync(img)

    words: list[dict] = []
    line_texts: list[str] = []
    for line in _winrt_iter(getattr(result, "lines", None)):
        ltext = (getattr(line, "text", "") or "").strip()
        if ltext:
            line_texts.append(ltext)
        for w in _winrt_iter(getattr(line, "words", None)):
            try:
                r = w.bounding_rect
                words.append({"text": w.text, "x0": r.x * scale, "x1": (r.x + r.width) * scale, "top": r.y * scale, "bottom": (r.y + r.height) * scale})
            except Exception:
                continue
    if not line_texts:
        flat = getattr(result, "text", "") or ""
        line_texts = [ln.strip() for ln in flat.splitlines() if ln.strip()]
    return words, line_texts


def _ocr_page(doc, page_index: int, dpi: int = 300):
    """Scanned statements (image-only PDFs with no text layer) can't be read by
    any text method, so OCR them. Tesseract (bundled, offline) is the primary
    engine; winocr is a last-resort fallback. Returns (words, line_texts): word
    boxes for table reconstruction plus each line's plain text as a guaranteed
    fallback so a scanned page never comes out empty. Returns (None, None) only
    when no OCR engine is reachable at all, so the caller can flag the page."""
    scale = 72.0 / dpi  # OCR boxes are in image pixels; convert back to PDF points
    if ocr.available():
        import io

        from PIL import Image

        page = doc[page_index]
        pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        try:
            return ocr.ocr_words(img, scale)
        except Exception:
            pass  # fall through to winocr
    return _winocr_page(doc, page_index, dpi, scale)


def pdf_to_excel(path: str, save_path: str) -> dict:
    """Turn each page's table into a worksheet, choosing the best strategy per
    page: gridline-ruled tables first (one row per transaction even when text
    wraps), then a positional word-grid for borderless text tables, and a
    Windows OCR fallback for scanned image-only pages. Falls back to
    pdfplumber's default detection if the smarter paths find nothing, so
    simple PDFs that already converted keep working."""
    import pdfplumber

    wb = Workbook()
    wb.remove(wb.active)
    table_count = 0
    scanned_pages: list[int] = []

    doc = fitz.open(path)
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                words = page.extract_words()
                if not words:
                    # No text layer on this page: it's a scan. Try OCR.
                    ocr_words, ocr_lines = _ocr_page(doc, i)
                    if ocr_words is None and ocr_lines is None:
                        scanned_pages.append(i + 1)  # OCR unavailable (not on Windows)
                        continue
                    rows = _grid_from_words(ocr_words, page.width) if ocr_words else []
                    if len(rows) < 2 and ocr_lines:
                        # The positional grid came out empty or too thin for
                        # this scan; fall back to the raw OCR text so the page
                        # is never blank, split into rough columns where the
                        # spacing allows. Columns may need light manual cleanup.
                        rows = [_split_line_cols(t) for t in ocr_lines]
                    rows = _trim_empty_columns(rows)
                    rows = _split_leading_date(rows)  # separate Date from Description on statements
                    if not rows:
                        scanned_pages.append(i + 1)
                        continue
                else:
                    rows = _ruled_table(page)
                    if rows is None:
                        rows = _grid_from_words(
                            [
                                {"x0": w["x0"], "x1": w["x1"], "top": w["top"], "bottom": w["bottom"], "text": w["text"]}
                                for w in words
                            ],
                            page.width,
                        )
                    if not rows:
                        default = page.extract_tables()
                        rows = [r for t in default for r in t if any((c or "").strip() for c in r)]
                if not rows:
                    continue
                table_count += 1
                ws = wb.create_sheet(title=f"Page{i + 1}"[:31])
                for row in rows:
                    ws.append(["" if c is None else c for c in row])
    finally:
        doc.close()

    if not wb.sheetnames:
        note = wb.create_sheet(title="Sheet1")
        if scanned_pages:
            note["A1"] = (
                "This PDF is a scanned image with no text layer, and the OCR engine "
                "could not be reached, so no text could be extracted from it."
            )
        else:
            note["A1"] = "No tables were detected in this PDF."
    wb.save(save_path)
    return {"tables_found": table_count, "scanned_pages": scanned_pages}


def _visual_lines(block):
    """PyMuPDF's own "lines" grouping is driven by the PDF's text-showing
    operators, not by what's visually on the same row. PDFs that position
    every word with its own explicit placement (PowerPoint's PDF export does
    this routinely) come out of get_text("dict") as one "line" per *word*,
    same y0/y1, but each its own entry. Re-cluster spans by vertical bbox
    overlap instead, so a sentence PyMuPDF fragmented into eight one-word
    lines becomes the single flowing line it actually is on the page."""
    raw = [ln for ln in block.get("lines", []) if ln.get("spans")]
    clusters = []
    for ln in raw:
        y0, y1 = ln["bbox"][1], ln["bbox"][3]
        target = None
        for c in clusters:
            overlap = min(y1, c["y1"]) - max(y0, c["y0"])
            if overlap > 0.5 * min(y1 - y0, c["y1"] - c["y0"]):
                target = c
                break
        if target is None:
            clusters.append({"y0": y0, "y1": y1, "spans": list(ln["spans"])})
        else:
            target["spans"].extend(ln["spans"])
            target["y0"] = min(target["y0"], y0)
            target["y1"] = max(target["y1"], y1)
    for c in clusters:
        c["spans"].sort(key=lambda s: s["bbox"][0])  # left-to-right reading order
    clusters.sort(key=lambda c: c["y0"])
    return clusters


def _add_text_block(slide, block, page_height_pt):
    lines = _visual_lines(block)
    if not lines:
        return
    x0, y0, x1, y1 = block["bbox"]
    # PDF text bboxes hug the glyphs tightly; give the textbox a little
    # breathing room so descenders/ascenders don't get clipped in PowerPoint.
    pad_pt = 2
    box = slide.shapes.add_textbox(
        Emu(int((x0 - pad_pt) * EMU_PER_POINT)),
        Emu(int((y0 - pad_pt) * EMU_PER_POINT)),
        Emu(int((x1 - x0 + 2 * pad_pt) * EMU_PER_POINT)),
        Emu(int((y1 - y0 + 2 * pad_pt) * EMU_PER_POINT)),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        spans = line["spans"]
        for j, span in enumerate(spans):
            text = _sanitize_text(span.get("text", ""))
            if not text:
                continue
            # Independently-placed words carry no space of their own, infer
            # one from the real gap between this word and the previous.
            if j > 0:
                gap = span["bbox"][0] - spans[j - 1]["bbox"][2]
                if gap > 1.0:
                    text = " " + text
            run = para.add_run()
            run.text = text
            size = span.get("size", 12)
            run.font.size = Pt(max(1, round(size)))
            flags = span.get("flags", 0)
            run.font.bold = bool(flags & PDF_BOLD_FLAG)
            run.font.italic = bool(flags & PDF_ITALIC_FLAG)
            color_int = span.get("color", 0)
            r = (color_int >> 16) & 255
            g = (color_int >> 8) & 255
            b = color_int & 255
            run.font.color.rgb = RGBColor(r, g, b)
            # Font *family* isn't preserved, PDF embeds subset fonts under
            # internal names (e.g. "ABCDEE+Calibri-Bold") that rarely match
            # an installed PowerPoint font, so guessing one in would more
            # often replace a correct-looking font with a wrong one.


def _add_page_images(slide, page, doc):
    for xref, *_ in page.get_images(full=True):
        rects = page.get_image_rects(xref)
        if not rects:
            continue
        rect = rects[0]
        try:
            info = doc.extract_image(xref)
        except Exception:
            continue
        ext = info.get("ext", "png")
        tmp = tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False)
        try:
            tmp.write(info["image"])
            tmp.close()
            slide.shapes.add_picture(
                tmp.name,
                Emu(int(rect.x0 * EMU_PER_POINT)),
                Emu(int(rect.y0 * EMU_PER_POINT)),
                Emu(int(rect.width * EMU_PER_POINT)),
                Emu(int(rect.height * EMU_PER_POINT)),
            )
        finally:
            Path(tmp.name).unlink(missing_ok=True)


def pdf_to_ppt(path: str, save_path: str) -> None:
    """Rebuilds each page as real, editable PowerPoint content: text blocks
    become editable text boxes (position, size, bold/italic, and color
    preserved; font family is not, see _add_text_block), and embedded
    images become separate picture shapes rather than one flattened
    page-sized raster. Works best on simple, single-column layouts; dense
    multi-column pages may need manual tidying after conversion, the same
    trade-off as PDF to Word."""
    doc = fitz.open(path)
    if doc.page_count == 0:
        doc.close()
        raise ValueError("This PDF has no pages.")

    first = doc[0]
    prs = Presentation()
    prs.slide_width = Emu(int(first.rect.width * EMU_PER_POINT))
    prs.slide_height = Emu(int(first.rect.height * EMU_PER_POINT))
    blank_layout = prs.slide_layouts[6]

    for page in doc:
        slide = prs.slides.add_slide(blank_layout)
        _add_page_images(slide, page, doc)
        text_dict = page.get_text("dict")
        for block in text_dict.get("blocks", []):
            if block.get("type") == 0:  # 0 = text, 1 = image (images handled separately above)
                _add_text_block(slide, block, first.rect.height)

    doc.close()
    prs.save(save_path)


def pdf_to_images(path: str, save_dir: str, fmt: str = "png", dpi: int = 150) -> list[str]:
    doc = fitz.open(path)
    base = Path(path).stem
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    outputs = []
    for page in doc:
        pix = page.get_pixmap(matrix=mat)
        out_path = str(Path(save_dir) / f"{base}-page{page.number + 1}.{fmt}")
        pix.save(out_path)
        outputs.append(out_path)
    doc.close()
    return outputs
